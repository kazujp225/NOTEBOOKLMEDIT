"""
Export Service - Generate PDF and PPTX from corrected pages
"""
import io
from typing import List
from uuid import UUID
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from backend.storage import storage
from backend.config import settings


def create_pdf(page_paths: List[str], output_path: str) -> str:
    """
    Create PDF from page images

    Args:
        page_paths: List of storage paths to page images (in order)
        output_path: Storage path for output PDF

    Returns:
        Storage path to created PDF
    """
    buffer = io.BytesIO()

    c = canvas.Canvas(buffer)

    for i, path in enumerate(page_paths):
        # Load image
        img_bytes = storage().get(path)
        img = Image.open(io.BytesIO(img_bytes))

        # Set page size to match image (in points, 72 dpi)
        # Convert from pixels at 300dpi to points at 72dpi
        width_pt = img.width * 72 / settings.pdf_dpi
        height_pt = img.height * 72 / settings.pdf_dpi

        c.setPageSize((width_pt, height_pt))

        # Draw image
        img_reader = ImageReader(io.BytesIO(img_bytes))
        c.drawImage(img_reader, 0, 0, width=width_pt, height=height_pt)

        c.showPage()

    c.save()

    # Save to storage
    storage().save_bytes(buffer.getvalue(), output_path)

    return output_path


def create_pptx(page_paths: List[str], output_path: str) -> str:
    """
    Create PPTX with each page as an image slide

    Args:
        page_paths: List of storage paths to page images (in order)
        output_path: Storage path for output PPTX

    Returns:
        Storage path to created PPTX
    """
    prs = Presentation()

    # Set slide size to widescreen (16:9) - can adjust
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Blank layout
    blank_layout = prs.slide_layouts[6]

    for i, path in enumerate(page_paths):
        # Load image
        img_bytes = storage().get(path)
        img = Image.open(io.BytesIO(img_bytes))

        # Add slide
        slide = prs.slides.add_slide(blank_layout)

        # Calculate dimensions to fit slide while maintaining aspect ratio
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        img_ratio = img.width / img.height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            # Image is wider - fit to width
            width = slide_width
            height = width / img_ratio
        else:
            # Image is taller - fit to height
            height = slide_height
            width = height * img_ratio

        # Center on slide
        left = (slide_width - width) / 2
        top = (slide_height - height) / 2

        # Save temp image for python-pptx
        img_buffer = io.BytesIO(img_bytes)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            left,
            top,
            width,
            height
        )

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)

    # Save to storage
    storage().save_bytes(buffer.getvalue(), output_path)

    return output_path


def get_page_paths_for_project(project_id: UUID, db) -> List[str]:
    """
    Get ordered list of page image paths for a project

    Args:
        project_id: Project UUID
        db: Database session

    Returns:
        List of storage paths in page order
    """
    from backend.db.models import Page

    pages = db.query(Page).filter(
        Page.project_id == project_id
    ).order_by(Page.page_number).all()

    return [page.image_path for page in pages]


def export_project_pdf(project_id: UUID, db) -> str:
    """
    Export project as PDF

    Args:
        project_id: Project UUID
        db: Database session

    Returns:
        Storage path to exported PDF
    """
    page_paths = get_page_paths_for_project(project_id, db)
    output_path = f"projects/{project_id}/exports/output.pdf"

    return create_pdf(page_paths, output_path)


def export_project_pptx(project_id: UUID, db) -> str:
    """
    Export project as PPTX

    Args:
        project_id: Project UUID
        db: Database session

    Returns:
        Storage path to exported PPTX
    """
    page_paths = get_page_paths_for_project(project_id, db)
    output_path = f"projects/{project_id}/exports/output.pptx"

    return create_pptx(page_paths, output_path)
